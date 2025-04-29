import json
import os
import tempfile
import uuid

from dotenv import load_dotenv
from fastapi import FastAPI, Request, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from typing import List
from openai import OpenAI
import PyPDF2
import docx
from pptx import Presentation

# load env variables
load_dotenv()

app = FastAPI()

# Configure CORS - using a more permissive configuration for troubleshooting
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # allow all origins (TODO unsafe)
    allow_credentials=False,  # must be False when using allow_origins=["*"]
    allow_methods=["*"],
    allow_headers=["*"],
)

# initialize OpenAI client with API key (export OPENAI_API_KEY="...")
client = OpenAI(api_key="YOUR_OPENAI_API_KEY") # TODO: replace with yours


@app.get("/test")
async def test_endpoint():
    return {"message": "Backend is working!"}


PROMPT = """You are an expert educational assistant. Your task is to help a user learn a specific topic by generating a list of up to 10 concise learning notes, each formatted as:
- Title: 1 short, specific line
- Content: A short paragraph or bullet points (ideally < 200 characters) explaining the most important idea. You may use simple Markdown formatting if helpful.

Each note should focus on one key idea, suitable for placement in a memory palace. To make learning more effective and engaging, vary the type of information across cards. Some possible types include:
- High-level explanation
- Core definition
- Why it matters
- Real-world application
- Key components
- Examples or analogies
- Common mistakes or misunderstandings
- Comparison with a related concept
- Historical background (if relevant)
- Fun facts

These are examples of card types. You can use different types depending on the topic. Don't include all of these types, they are just examples. Focus on clarity, usefulness, and learning value. The cards should be varied and useful to the learner.

For instance, if I ask about neural networks, I would like a basic definition of neural networks, of common activation functions and why they're used, why neural networks are useful and important, different popular types of neural networks (MLP, LSTM, CNN, Transformers...), challenges with overfitting/underfitting, examples of achievements using neural networks, and so on. Try to create a cohesive story and give useful information (not just superficial facts)."""


def query_gpt(user_prompt):
    print("REQUEST WITH PROMPT", user_prompt)
    if len(user_prompt) > 100:
        print(f"User prompt too long ({len(user_prompt)} characters). Truncated to 100 characters.")
        user_prompt = user_prompt[:100]
    user_prompt = "I want to learn about " + user_prompt

    # model: gpt-4o, $2.50/1M tokens (debugging: gpt-4o-mini, $0.15/1M tokens)
    response = client.responses.create(
        model="gpt-4o-mini",
        input=[
            {"role": "developer", "content": PROMPT},
            {"role": "user", "content": user_prompt},
        ],
        text={
            "format": {
                "type": "json_schema",
                "name": "calendar_event",
                "schema": {
                    "type": "object",
                    "properties": {
                        "wrapper": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "title": {"type": "string"},
                                    "content": {"type": "string"},
                                },
                                "required": ["title", "content"],
                                "additionalProperties": False,
                            },
                        },
                    },
                    "required": ["wrapper"],
                    "additionalProperties": False,
                },
                "strict": True,
            }
        },
    )
    print("REQUEST ENDED")

    output = json.loads(response.output_text)
    cards = output["wrapper"]

    for card in cards:
        yield (card["title"], card["content"])


last_id = 0


@app.post("/generate-notes")
async def generate_notes(request: Request):
    global last_id
    data = await request.json()
    user_prompt = data.get("prompt", "")

    cards = []
    for title, content in query_gpt(user_prompt):
        cards.append(
            {
                "id": last_id,
                "title": title,
                "content": content,
            }
        )
        last_id += 1

    return cards


@app.post("/upload-files")
async def upload_files(files: List[UploadFile] = File(...)):
    global last_id
    cards = []
    
    for file in files:
        file_content = await process_file(file)
        if file_content:
            # Create a prompt based on the file content
            file_prompt = f"Generate notes from this content: {file_content[:1000]}..."
            
            # Process the content with GPT
            for title, content in query_gpt(file_prompt):
                cards.append(
                    {
                        "id": last_id,
                        "title": title,
                        "content": content,
                    }
                )
                last_id += 1
    
    return cards

async def process_file(file: UploadFile):
    try:
        # Create a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.filename)[1]) as temp_file:
            # Write uploaded file content to the temp file
            content = await file.read()
            temp_file.write(content)
            temp_file_path = temp_file.name
        
        # Extract text based on file type
        file_extension = os.path.splitext(file.filename)[1].lower()
        extracted_text = ""
        
        if file_extension == '.pdf':
            # Process PDF
            with open(temp_file_path, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    extracted_text += page.extract_text() + "\n"
        
        elif file_extension in ['.docx', '.doc']:
            # Process Word document
            doc = docx.Document(temp_file_path)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"
        
        elif file_extension in ['.pptx', '.ppt']:
            # Process PowerPoint
            prs = Presentation(temp_file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
        
        elif file_extension in ['.mp4', '.mov', '.avi']:
            # For video files, just acknowledge we can't process them directly
            # In a real implementation, you might use a transcription service
            extracted_text = f"Content from video file: {file.filename}. Video content requires transcription services."
        
        else:
            # For other file types, read as text if possible
            extracted_text = content.decode('utf-8', errors='ignore')
        
        # Clean up temp file
        os.unlink(temp_file_path)
        
        return extracted_text
    
    except Exception as e:
        print(f"Error processing file {file.filename}: {str(e)}")
        return f"Error processing file {file.filename}. Please try a different file format."


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000, reload=True)
